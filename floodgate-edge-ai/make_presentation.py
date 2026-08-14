#!/usr/bin/env python3
"""
make_presentation.py — generates the FloodGate Edge Brain demo deck
(FloodGate_Edge_Brain_Presentation.pptx, 16:9, dark theme).

Run:  python3 make_presentation.py        (uses the workspace venv: .venv)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
BG      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
PANEL   = RGBColor(0x1B, 0x2B, 0x42)   # panel navy
PANEL2  = RGBColor(0x24, 0x3B, 0x5A)   # lighter panel
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)   # accent
AMBER   = RGBColor(0xFB, 0xBF, 0x24)   # highlight
RED     = RGBColor(0xF8, 0x71, 0x71)   # warn
GREEN   = RGBColor(0x4A, 0xDE, 0x80)   # agree
WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
DARKTXT = RGBColor(0x0B, 0x12, 0x1E)

FONT = "Calibri"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def add_slide(prs, kicker, title, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # kicker
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = kicker.upper()
    r.font.size = Pt(12); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = TEAL
    # title
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = WHITE
    # accent underline
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.52),
                                Inches(1.6), Pt(3.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = TEAL; ln.line.fill.background()
    ln.shadow.inherit = False
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def footer(slide, n):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "FloodGate Edge Brain"
    r.font.size = Pt(9); r.font.name = FONT; r.font.color.rgb = MUTED
    r2 = p.add_run(); r2.text = f"   ·   {n}"
    r2.font.size = Pt(9); r2.font.name = FONT; r2.font.color.rgb = MUTED


def box(slide, text, left, top, w, h, fill=PANEL, color=WHITE, size=14,
        bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.LEFT,
        mono=False):
    sp = slide.shapes.add_shape(shape, left, top, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = RGBColor(0x33, 0x4A, 0x6E); sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = MONO if mono else FONT
    r.font.color.rgb = color
    return sp


def bullets(slide, items, left, top, width, height, size=15, gap=6):
    """items: list of (level, text) or (level, text, color)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        level = it[0]
        text = it[1]
        color = it[2] if len(it) > 2 else WHITE
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        mark = "▸  " if level == 0 else "–  "
        r = p.add_run(); r.text = mark
        r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = TEAL if level == 0 else MUTED
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(size); r2.font.name = FONT; r2.font.color.rgb = color
    return tb


def chip_row(slide, labels, left, top, w, h, fill=PANEL2, color=WHITE, size=13,
             gap=0.15, accent=TEAL):
    x = left
    for label in labels:
        box(slide, label, x, top, w, h, fill=fill, color=color, size=size,
            bold=True, align=PP_ALIGN.CENTER)
        # tiny accent dot
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.08), top + h // 2 - Pt(3),
                                   Pt(6), Pt(6))
        d.fill.solid(); d.fill.fore_color.rgb = accent; d.line.fill.background()
        d.shadow.inherit = False
        x += w + gap
    return x


def arrow(slide, left, top, w=Inches(0.5), h=Inches(0.3), color=TEAL):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    a.shadow.inherit = False
    return a


def table(slide, rows, cols, left, top, width, height, header_fill=TEAL,
          header_color=DARKTXT, cell_fill=PANEL, cell_color=WHITE, size=13,
          col_widths=None):
    gfx = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else cell_fill
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    return tbl


def fill_cell(tbl, r, c, text, color=WHITE, bold=False, size=13):
    cell = tbl.cell(r, c)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.name = FONT
    run.font.color.rgb = color


# ============================================================================
# Deck
# ============================================================================
prs = new_deck()

# ---------------------------------------------------------------- S1 title
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
bg.shadow.inherit = False
box(s, "FLOODGATE  ·  SMART STORM DRAIN MONITORING  ·  TEAM DEMO ADD-ON",
    Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.45), fill=PANEL,
    color=TEAL, size=13, bold=True, align=PP_ALIGN.CENTER)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.6))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "FloodGate Edge Brain"
r.font.size = Pt(54); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
tb2 = s.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.0))
tf = tb2.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "An on-device AI decision layer for flood detection — a 14 MB language model "
r.font.size = Pt(20); r.font.name = FONT; r.font.color.rgb = MUTED
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "running on a separate Raspberry Pi, fully offline"
r.font.size = Pt(20); r.font.name = FONT; r.font.color.rgb = MUTED
stats = ["14 MB model", "45M parameters", "< 0.5 s per decision", "100% offline"]
chip_row(s, stats, Inches(1.7), Inches(5.0), Inches(2.3), Inches(0.5),
         fill=PANEL2, color=WHITE, size=13, gap=0.25, accent=AMBER)
tb = s.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Add-on on its own Pi  —  production gateway untouched  ·  Needle 2 (Cactus Compute)  ·  LoRA fine-tuned"
r.font.size = Pt(12); r.font.name = FONT; r.font.color.rgb = MUTED
s.notes_slide.notes_text_frame.text = (
    "Good morning. Today we show the FloodGate Edge Brain: a tiny AI model that "
    "lives on its own Raspberry Pi, watches the same sensor stream as the rest of "
    "the system, and decides what alert to raise — fully offline, in well under a second.")

# ---------------------------------------------------------------- S2 what
s = add_slide(prs, "Overview", "What is the Edge Brain?",
              notes="One sentence: a local AI decision layer that reads the sensor stream "
                    "and picks the alert, while the rule engine keeps the final word.")
bullets(s, [
    (0, "A local AI decision layer that watches the SAME MQTT sensor stream as the "
        "existing bridge and dashboard — and every cycle decides what alert to raise."),
    (0, "Powered by Needle 2: a 45M-parameter tool-calling LLM, small enough to run "
        "on a Raspberry Pi 4/5 with no GPU and no internet."),
    (0, "The deterministic rule engine still drives the real alert — the AI is "
        "compared against it every cycle (AGREE / DISAGREE), never blindly trusted."),
    (0, "Runs as its own service on a separate Pi (:8090) — zero changes to the "
        "production gateway that runs the bridge and dashboard."),
], Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.2), size=17, gap=14)
box(s, "Sensor → decision in well under a second, 100% on-device — "
       "for flood infrastructure, offline-first is the safety requirement.",
    Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.9), fill=PANEL, color=TEAL,
    size=16, bold=True, align=PP_ALIGN.CENTER)
footer(s, 2)

# ---------------------------------------------------------------- S3 why
s = add_slide(prs, "Why", "Why an AI at the edge?",
              notes="The honest pitch: rules are the better tool for thresholds; the AI "
                    "earns its place around them — implausibility, explanations, doubt.")
bullets(s, [
    (0, "Offline-first safety — real floods kill connectivity; every decision stays local."),
    (0, "Private, instant, free — no per-call cost, no sensor data leaves the device."),
    (0, "Rules only cover what someone enumerated — combinations explode; a model "
        "learns the shape of the decision and interpolates."),
    (0, "Rules can't recognize “this input makes no sense” — frozen readings, a 50 cm "
        "jump in 20 seconds, a dead ADC; the model learns the concept of implausibility."),
    (0, "Rules output codes; the model outputs the WHY — the explanation is half the "
        "product of a public warning system."),
    (0, "Rules have cliff edges — the model can express doubt, which drives escalation."),
], Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.4), size=17, gap=12)
footer(s, 3)

# ---------------------------------------------------------------- S4 architecture
s = add_slide(prs, "Architecture", "Where it fits in FloodGate",
              notes="Same MQTT bus, three listeners. The edge brain is an add-on on a "
                    "separate Pi; it only reads the bus (plus optional cloud escalation).")
# ESP32
box(s, "ESP32 edge node\nBMP280 + water depth\npublishes every 1–5 min (deep sleep)",
    Inches(0.6), Inches(2.4), Inches(2.9), Inches(1.5), fill=PANEL2, size=13)
# broker
box(s, "MQTT broker\numd/cpse/floodgate/telemetry", Inches(3.9), Inches(2.7),
    Inches(2.4), Inches(0.9), fill=PANEL, size=13, bold=True)
arrow(s, Inches(3.55), Inches(2.95), Inches(0.3), Inches(0.25))
arrow(s, Inches(6.35), Inches(2.95), Inches(0.3), Inches(0.25))
# main gateway
box(s, "MAIN GATEWAY PI 4 (unchanged)\nbridge → Supabase   ·   dashboard :8080",
    Inches(7.0), Inches(1.9), Inches(5.6), Inches(1.2), fill=PANEL, size=13)
# edge brain pi
box(s, "EDGE BRAIN PI (separate, this demo)\nsubscribes to the same topic   ·   API :8090",
    Inches(7.0), Inches(3.5), Inches(5.6), Inches(1.2), fill=RGBColor(0x13, 0x3A, 0x3A),
    size=13, bold=True)
# cloud
box(s, "cloud tier (optional)\nonly when unsure", Inches(7.0), Inches(5.1),
    Inches(5.6), Inches(0.9), fill=PANEL, size=12)
# weather note
box(s, "OpenWeather “high risk” broadcasts ride the same topic → the brain hears them too",
    Inches(0.6), Inches(4.5), Inches(5.7), Inches(0.9), fill=PANEL, size=12, color=MUTED)
# takeaway
box(s, "Add-on design: the brain only LISTENS to the existing bus — no firmware, bridge or "
       "dashboard changes; it can even be demoed with a simulator on a Pi with no sensors.",
    Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.85), fill=PANEL, color=AMBER,
    size=13, bold=True)
footer(s, 4)

# ---------------------------------------------------------------- S5 layers
s = add_slide(prs, "Design", "Three layers, one decision",
              notes="The transferable pattern: rules for what must never fail, a small local "
                    "model for what can't be enumerated, a cloud model for what's rare and hard.")
box(s, "L3  Cloud model (Anthropic, optional)\nonly the rare, genuinely hard cases — "
       "uncertainty, conflict, sensor faults",
    Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.2), fill=PANEL, size=15)
box(s, "L2  Small local model — Needle 2 (45M)\npattern & plausibility, explanations, "
       "doubt · compared to the rules every cycle (AGREE/DISAGREE)",
    Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.2), fill=PANEL2, size=15)
box(s, "L1  Deterministic rule engine — always drives the real alert\nauditable, instant, "
       "can never hallucinate  →  the AI is never a single point of failure",
    Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.2), fill=TEAL, color=DARKTXT,
    size=15, bold=True)
footer(s, 5)

# ---------------------------------------------------------------- S6 contract
s = add_slide(prs, "Data contract", "One wire format, designed once",
              notes="The whole system shares one payload. Units convert once at the boundary; "
                    "the sensor reports its own health.")
box(s, '{"device_id": "ESP32-FloodGate",\n'
       ' "water_depth": 0.245,          // METERS\n'
       ' "atm_pressure_hpa": 1013.2,\n'
       ' "ambient_temp_c": 24.1,\n'
       ' "status": "OK"}                // "OK" | "ADC_ERROR"',
    Inches(0.7), Inches(1.95), Inches(5.6), Inches(2.5), fill=RGBColor(0x0B, 0x12, 0x1E),
    color=GREEN, size=13, mono=True)
bullets(s, [
    (0, "Units convert at the boundary, once — meters → cm in a single parser; "
        "everything downstream (rules, prompts, training) thinks in cm only."),
    (0, "The sensor tells you when it is lying — status ADC_ERROR and derived "
        "STALE (no reading for 15 min) are first-class features."),
    (0, "Bad status escalates BEFORE any water rule can fire — we never raise a "
        "flood alert on data we know is garbage."),
    (0, "fg_core.py is the single source of truth — thresholds, parser, rules and "
        "the prompt format; training, runtime and eval all import it."),
], Inches(0.7), Inches(4.75), Inches(11.9), Inches(2.1), size=15, gap=10)
footer(s, 6)

# ---------------------------------------------------------------- S7 rules
s = add_slide(prs, "Safety net", "The rule engine — ten lines, always in charge",
              notes="This table is the ground truth the model is graded against and the "
                    "alert authority in production.")
tbl = table(s, 6, 2, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.6),
            col_widths=[Inches(7.6), Inches(4.3)])
fill_cell(tbl, 0, 0, "Condition", bold=True)
fill_cell(tbl, 0, 1, "Action", bold=True)
rules = [
    ("sensor status ≠ OK  (ADC_ERROR / STALE)", "escalate_to_cloud"),
    ("water level ≥ 35 cm", "issue_flood_warning"),
    ("water level ≥ 25 cm  or  rise ≥ 0.5 cm/min", "issue_flood_watch"),
    ("pressure trend ≤ −2 hPa/hr  or  weather high-risk", "issue_storm_watch"),
    ("otherwise", "report_all_clear"),
]
for i, (c, a) in enumerate(rules, start=1):
    fill_cell(tbl, i, 0, c)
    fill_cell(tbl, i, 1, a, color=TEAL, bold=True)
box(s, "Auditable · instant · deterministic  —  a judge can point at the exact line that fired",
    Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.7), fill=PANEL, color=WHITE,
    size=14, bold=False, align=PP_ALIGN.CENTER)
footer(s, 7)

# ---------------------------------------------------------------- S8 tools
s = add_slide(prs, "The AI", "Five tools the model can call",
              notes="The tool names are the alert vocabulary shared with the dashboard. "
                    "The docstrings are the model's instructions.")
labels = ["report_all_clear", "issue_storm_watch", "issue_flood_watch",
          "issue_flood_warning", "escalate_to_cloud"]
chip_row(s, labels, Inches(0.7), Inches(1.9), Inches(2.35), Inches(0.55),
         fill=PANEL2, size=12, gap=0.1, accent=TEAL)
box(s, '[14:02:31] edge=report_all_clear     rule=report_all_clear     AGREE    42ms\n'
       '[14:03:41] edge=issue_storm_watch    rule=issue_storm_watch    AGREE    38ms\n'
       '[14:05:03] edge=escalate_to_cloud    rule=escalate_to_cloud    AGREE    41ms',
    Inches(0.7), Inches(2.85), Inches(7.0), Inches(2.0),
    fill=RGBColor(0x0B, 0x12, 0x1E), color=GREEN, size=13, mono=True)
bullets(s, [
    (0, "Every cycle the brain asks the model to pick exactly one tool for the "
        "current sensor summary."),
    (0, "The tool docstrings ARE the model's instructions — it reads them with "
        "every prompt."),
    (0, "Its answer is compared to the rule engine; AGREE/DISAGREE is printed live "
        "and pushed to the dashboard."),
], Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.8), size=15, gap=8)
footer(s, 8)

# ---------------------------------------------------------------- S9 escalation
s = add_slide(prs, "Edge-cloud", "When the edge asks for a second opinion",
              notes="Routine decisions stay local; only the rare hard cases pay the cloud cost.")
bullets(s, [
    (0, "The system escalates to a cloud model when:"),
    (1, "the edge model explicitly calls escalate_to_cloud,"),
    (1, "edge and rule DISAGREE,"),
    (1, "model confidence drops below the floor (0.6),"),
    (1, "the sensor status is faulty or stale."),
    (0, "Routine decisions stay local, private, instant and free; only the rare, "
        "genuinely hard cases use a frontier model."),
    (0, "Without an API key, escalations are logged and skipped — the demo still "
        "works 100% offline."),
], Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.4), size=16, gap=12)
box(s, "Small edge model = intelligent filter for a large cloud model — "
       "one of the most reusable ideas in this project",
    Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.9), fill=PANEL, color=AMBER,
    size=15, bold=True, align=PP_ALIGN.CENTER)
footer(s, 9)

# ---------------------------------------------------------------- S10 training
s = add_slide(prs, "Training", "From 475 flashcards to a 23 MB deployable",
              notes="Dataset is derived from the rules + adversarial fault cases; LoRA on a "
                    "laptop; locked toolchain; 4-bit quantization.")
steps = ["make_finetune_data.py\n475 flashcards\n(rule-consistent + fault cases)",
         "finetune_fp32.py\nLoRA rank 16, 4 epochs\nloss 2.5 → 1.50",
         "needle build --bits 4\nmerge + quantize\nfloodgate.cact (23 MB)"]
x = Inches(0.6)
for i, t in enumerate(steps):
    box(s, t, x, Inches(2.0), Inches(3.55), Inches(1.7), fill=PANEL2, size=13)
    if i < 2:
        arrow(s, x + Inches(3.6), Inches(2.6), Inches(0.4), Inches(0.3))
    x += Inches(4.0)
bullets(s, [
    (0, "The dataset is derived from the same rules the model is graded against — "
        "plus deliberate adversarial cases (frozen / conflicting / spiking sensors) "
        "where the model must OVERRULE the rules."),
    (0, "fg_core.py guarantees training prompts and live prompts are byte-identical."),
    (0, "Locked toolchain (cactus-needle 2.0.2 + pinned engine): a silent upstream "
        "change in the trainer's target layers made everything worse — reproducibility "
        "beat “latest”."),
    (0, "2-bit quantization broke the model; 4-bit is the sweet spot for exact JSON "
        "tool calls."),
], Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.6), size=15, gap=9)
footer(s, 10)

# ---------------------------------------------------------------- S11 results
s = add_slide(prs, "Results", "Measured, not claimed",
              notes="200 held-out scenarios, rule engine as ground truth. Honest numbers: "
                    "the model masters calm/pressure; water thresholds stay with the rules.")
tbl = table(s, 6, 3, Inches(0.7), Inches(1.9), Inches(8.6), Inches(3.3),
            col_widths=[Inches(3.9), Inches(2.35), Inches(2.35)])
fill_cell(tbl, 0, 0, "Class (held-out, n=40 each)", bold=True)
fill_cell(tbl, 0, 1, "Base", bold=True)
fill_cell(tbl, 0, 2, "Tuned", bold=True)
res = [
    ("report_all_clear", "0%", "85%"),
    ("issue_storm_watch", "80%", "70%"),
    ("flood watch / warning / escalate", "0%", "0%"),
    ("OVERALL (n=200)", "16%", "31%"),
]
for i, (c, a, b) in enumerate(res, start=1):
    fill_cell(tbl, i, 0, c, bold=(i == 4))
    fill_cell(tbl, i, 1, a, bold=(i == 4))
    fill_cell(tbl, i, 2, b, color=GREEN if i < 3 else AMBER, bold=True)
bullets(s, [
    (0, "The 45M model reliably masters the calm & pressure classes — that's the "
        "visible AGREE in the demo."),
    (0, "Water thresholds stay with the rule engine — by design, not by accident; "
        "every disagreement escalates."),
    (0, "We pushed hard on the weak classes: 6 controlled runs (rebalanced data, "
        "class-weighted loss, rank 32) — the model can't even memorize them "
        "(train-fit 0–7%); flood_warning only improves at the cost of all-clear. "
        "That's the ceiling of a 45M model, verified, not assumed."),
    (0, "Honest limit: synthetic data mirrors the rules, so agreement partly "
        "reflects imitation — real storm data is the next step."),
], Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.8), size=13, gap=6)
footer(s, 11)

# ---------------------------------------------------------------- S12 robustness
s = add_slide(prs, "Engineering", "Four bugs we found, four fixes",
              notes="Each one is a real engineering lesson, all documented in the tutorial.")
war = [
    ("loss = NaN in the vendor trainer", "half-precision LoRA underflowed AdamW's eps → "
     "fp32 adapter + grad clipping (finetune_fp32.py)"),
    (".cact format version skew", "exporter wrote tag …82, engine read …83 → pinned "
     "2.0.0 engine + locked toolchain (pin_engine.py)"),
    ("2-bit quantization too lossy", "model collapsed into “reasoning loops” → build at "
     "4-bit, measured on our own task"),
    ("macOS engine segfaults mid-eval", "external weights crash after ~6–8 calls → "
     "crash-isolated inference worker (infer_worker.py)"),
]
y = Inches(1.95)
for title, body in war:
    box(s, title, Inches(0.7), y, Inches(3.6), Inches(1.15), fill=PANEL2, size=14, bold=True)
    box(s, body, Inches(4.5), y, Inches(8.1), Inches(1.15), fill=PANEL, size=12.5)
    y += Inches(1.3)
footer(s, 12)

# ---------------------------------------------------------------- S13 demo plan
s = add_slide(prs, "Demo", "Live demo on the separate Pi",
              notes="Simulator, four phases, ~45s each. Watch AGREE lines, rule alerts, "
                    "escalation, and the kill-WiFi moment.")
phases = [
    ("1 · Normal", "AGREE — all clear\nmodel learned calm conditions"),
    ("2 · Storm", "AGREE — storm watch\npressure drop + weather broadcast"),
    ("3 · Flood", "RULE alert + DISAGREE → escalation\n(model is weak on water thresholds)"),
    ("4 · Sensor fault", "both layers escalate\nstatus = ADC_ERROR"),
]
x = Inches(0.6)
for title, body in phases:
    box(s, title, x, Inches(2.0), Inches(2.85), Inches(0.6), fill=TEAL,
        color=DARKTXT, size=14, bold=True, align=PP_ALIGN.CENTER)
    box(s, body, x, Inches(2.75), Inches(2.85), Inches(1.6), fill=PANEL2, size=12.5)
    x += Inches(3.1)
bullets(s, [
    (0, "Decision line printed every ~10 s: edge=… rule=… AGREE/DISAGREE …ms."),
    (0, "Kill the WiFi → everything keeps working. That is the money moment."),
    (0, "Feed a conflicting reading → watch it escalate to the cloud tier."),
], Inches(0.6), Inches(4.75), Inches(12.1), Inches(1.9), size=15, gap=8)
footer(s, 13)

# ---------------------------------------------------------------- S14 deploy
s = add_slide(prs, "Deployment", "Mac trains, Pi runs — one command each",
              notes="Fully reproducible. The bundle is self-contained; deploy_pi.sh checks "
                    "for 64-bit OS, installs the service, evaluates on-device.")
box(s, "MAC  (train)", Inches(0.7), Inches(2.0), Inches(5.6), Inches(0.6),
    fill=PANEL2, size=15, bold=True, align=PP_ALIGN.CENTER)
box(s, "./setup_mac.sh", Inches(0.7), Inches(2.7), Inches(5.6), Inches(0.8),
    fill=RGBColor(0x0B, 0x12, 0x1E), color=GREEN, size=16, mono=True, align=PP_ALIGN.CENTER)
bullets(s, [
    (0, "env (locked toolchain) → 475-example dataset → LoRA train → 4-bit build → "
        "smoke test → bundle"),
], Inches(0.7), Inches(3.7), Inches(5.6), Inches(1.2), size=13)
box(s, "SEPARATE PI 4/5  (run, 64-bit OS)", Inches(6.9), Inches(2.0), Inches(5.7),
    Inches(0.6), fill=PANEL2, size=15, bold=True, align=PP_ALIGN.CENTER)
box(s, "tar xzf bundle && bash deploy_pi.sh", Inches(6.9), Inches(2.7), Inches(5.7),
    Inches(0.8), fill=RGBColor(0x0B, 0x12, 0x1E), color=GREEN, size=15, mono=True,
    align=PP_ALIGN.CENTER)
bullets(s, [
    (0, "venv (inference-only, no jax) → engine pin → on-device eval → systemd service"),
    (0, "status API at :8090 — the production gateway keeps :8080 untouched"),
], Inches(6.9), Inches(3.7), Inches(5.7), Inches(1.4), size=13)
box(s, "Docs: document/DESIGN.md (system design)  ·  document/TUTORIAL.md (walkthrough)",
    Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.8), fill=PANEL, color=TEAL,
    size=14, bold=True, align=PP_ALIGN.CENTER)
footer(s, 14)

# ---------------------------------------------------------------- S15 thanks
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
bg.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank you — questions?"
r.font.size = Pt(44); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
stats = ["14 MB model", "475 flashcards", "16% → 31% held-out", "4 war stories → 4 fixes",
         "1 separate Pi"]
chip_row(s, stats, Inches(1.2), Inches(4.4), Inches(2.1), Inches(0.55),
         fill=PANEL2, size=12, gap=0.18, accent=TEAL)
tb = s.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "floodgate-edge-ai · fg_core.py (single source of truth) · rules always in charge"
r.font.size = Pt(12); r.font.name = FONT; r.font.color.rgb = MUTED
s.notes_slide.notes_text_frame.text = (
    "Recap the three layers, the honest numbers, and invite questions about "
    "the war stories or the live demo.")

prs.save("FloodGate_Edge_Brain_Presentation.pptx")
print("saved FloodGate_Edge_Brain_Presentation.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
